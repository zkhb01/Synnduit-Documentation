"""Extract binary office docs (.docx/.pptx/.xlsx) and .rtf to Markdown next to the original.

Usage: python extract_docs.py [root]
Writes <original>.md beside each source file. Skips if up to date unless --force.
"""
import sys
import os
from pathlib import Path

ROOT = Path(sys.argv[1]) if len(sys.argv) > 1 and not sys.argv[1].startswith("--") else Path(".")
FORCE = "--force" in sys.argv


def md_escape(text: str) -> str:
    return text.replace("\r", "").strip()


def docx_to_md(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    lines = [f"# {path.stem}\n", f"> Extracted from `{path.name}`\n"]

    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.ns import qn

    def iter_block_items(document):
        # Iterate body children in document order; use the Document as the proxy
        # parent (it exposes .part, which the raw body element does not).
        body = document.element.body
        for child in body.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, document)
            elif child.tag == qn("w:tbl"):
                yield Table(child, document)

    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            txt = md_escape(block.text)
            if not txt:
                continue
            style = (block.style.name or "").lower() if block.style else ""
            if style.startswith("heading 1"):
                lines.append(f"\n## {txt}\n")
            elif style.startswith("heading 2"):
                lines.append(f"\n### {txt}\n")
            elif style.startswith("heading 3"):
                lines.append(f"\n#### {txt}\n")
            elif style.startswith("title"):
                lines.append(f"\n## {txt}\n")
            elif style.startswith("list") or block.text.lstrip().startswith(("-", "•", "*")):
                lines.append(f"- {txt.lstrip('-•* ')}")
            else:
                lines.append(f"{txt}\n")
        else:  # Table
            rows = block.rows
            if not rows:
                continue
            lines.append("")
            header = [md_escape(c.text).replace("\n", " ") or " " for c in rows[0].cells]
            lines.append("| " + " | ".join(header) + " |")
            lines.append("| " + " | ".join("---" for _ in header) + " |")
            for r in rows[1:]:
                cells = [md_escape(c.text).replace("\n", " ") or " " for c in r.cells]
                lines.append("| " + " | ".join(cells) + " |")
            lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def pptx_to_md(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    lines = [f"# {path.stem}\n", f"> Extracted from `{path.name}` ({len(prs.slides)} slides)\n"]
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## Slide {i}\n")
        # Title first if present
        title_txt = None
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title_txt = md_escape(slide.shapes.title.text)
            if title_txt:
                lines.append(f"**{title_txt}**\n")
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                tbl = shape.table
                rows = list(tbl.rows)
                if rows:
                    header = [md_escape(c.text).replace("\n", " ") or " " for c in rows[0].cells]
                    lines.append("| " + " | ".join(header) + " |")
                    lines.append("| " + " | ".join("---" for _ in header) + " |")
                    for r in rows[1:]:
                        cells = [md_escape(c.text).replace("\n", " ") or " " for c in r.cells]
                        lines.append("| " + " | ".join(cells) + " |")
                    lines.append("")
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = md_escape("".join(run.text for run in para.runs) or para.text)
                    if txt:
                        bullet = "  " * max(0, para.level) + "- "
                        lines.append(f"{bullet}{txt}")
        # Speaker notes
        if slide.has_notes_slide:
            notes = md_escape(slide.notes_slide.notes_text_frame.text)
            if notes:
                lines.append(f"\n*Notes:* {notes}\n")
    return "\n".join(lines).rstrip() + "\n"


def xlsx_to_md(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    lines = [f"# {path.stem}\n", f"> Extracted from `{path.name}`\n"]
    for ws in wb.worksheets:
        rows = [
            [("" if c is None else str(c)).replace("\n", " ").strip() for c in row]
            for row in ws.iter_rows(values_only=True)
        ]
        # Trim fully empty trailing rows/cols
        rows = [r for r in rows if any(cell for cell in r)]
        if not rows:
            continue
        # Trim trailing empty columns
        width = max(len(r) for r in rows)
        keep = [i for i in range(width) if any((r[i] if i < len(r) else "") for r in rows)]
        lines.append(f"\n## Sheet: {ws.title}\n")
        if not keep:
            continue
        def cell(r, i):
            return (r[i] if i < len(r) else "") or " "
        header = [cell(rows[0], i) for i in keep]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join("---" for _ in keep) + " |")
        for r in rows[1:]:
            lines.append("| " + " | ".join(cell(r, i) for i in keep) + " |")
        lines.append("")
    wb.close()
    return "\n".join(lines).rstrip() + "\n"


def rtf_to_md(path: Path) -> str:
    from striprtf.striprtf import rtf_to_text
    raw = path.read_text(encoding="utf-8", errors="replace")
    text = rtf_to_text(raw)
    lines = [f"# {path.stem}\n", f"> Extracted from `{path.name}`\n", ""]
    for ln in text.splitlines():
        s = ln.rstrip()
        lines.append(s)
    # collapse 3+ blank lines
    out = "\n".join(lines)
    while "\n\n\n\n" in out:
        out = out.replace("\n\n\n\n", "\n\n\n")
    return out.rstrip() + "\n"


HANDLERS = {".docx": docx_to_md, ".pptx": pptx_to_md, ".xlsx": xlsx_to_md, ".rtf": rtf_to_md}


def main():
    targets = []
    for p in ROOT.rglob("*"):
        if p.suffix.lower() in HANDLERS and not p.name.startswith("~$"):
            targets.append(p)
    targets.sort()
    print(f"Found {len(targets)} file(s) to extract under {ROOT.resolve()}\n")
    ok, skip, err = 0, 0, 0
    for src in targets:
        out = src.with_suffix(src.suffix + ".md")
        if out.exists() and not FORCE and out.stat().st_mtime >= src.stat().st_mtime:
            print(f"  SKIP (current)  {src}")
            skip += 1
            continue
        try:
            md = HANDLERS[src.suffix.lower()](src)
            out.write_text(md, encoding="utf-8")
            print(f"  OK  {src.name}  ->  {out.name}  ({len(md)} chars)")
            ok += 1
        except Exception as e:
            print(f"  ERROR  {src}: {type(e).__name__}: {e}")
            err += 1
    print(f"\nDone. {ok} written, {skip} skipped, {err} errors.")


if __name__ == "__main__":
    main()
